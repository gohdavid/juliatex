#!/usr/bin/env python3
"""Generate presentation slides from main.tex using the MIT PowerPoint template."""

import os, tempfile
from pptx import Presentation
from pptx.util import Pt, Inches
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
plt.rcParams['text.usetex'] = True
plt.rcParams['text.latex.preamble'] = r'\usepackage{amsmath,amssymb,bm}'
from PIL import Image

TEMPLATE = "MIT-PowerPoint-template-Arial.pptx"
OUTPUT = "presentation.pptx"
EQ_DIR = tempfile.mkdtemp(prefix="pptx_eq_")

prs = Presentation(TEMPLATE)

# ---------- delete existing slides ----------
for _ in range(len(prs.slides)):
    sldId = prs.slides._sldIdLst[0]
    rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    if rId is None:
        rId = sldId.get("r:id")
    if rId and rId in prs.part.rels:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sldId)

# ---------- layouts ----------
def get_layout(name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise ValueError(f"Layout '{name}' not found")

title_layout   = get_layout("Title Slide –\xa0Main, White, Red Text")
section_layout = get_layout("Section Divider –\xa0White, Red Text")
content_layout = get_layout("Title and Text")

# ---------- equation rendering ----------
_eq_counter = 0

def render_eq(latex_str, fontsize=20, max_width_in=7.0):
    """Render a LaTeX string (in math mode) to a transparent PNG.
    Returns (path, width_emu, height_emu)."""
    global _eq_counter
    _eq_counter += 1
    path = os.path.join(EQ_DIR, f"eq_{_eq_counter}.png")

    # Estimate figure height from number of lines
    nlines = latex_str.count(r'\\') + 1
    fig_h = max(0.8, 0.55 * nlines)
    fig, ax = plt.subplots(figsize=(10, fig_h))
    ax.text(0.0, 0.5, f"${latex_str}$",
            fontsize=fontsize, ha='left', va='center')
    ax.axis('off')
    fig.savefig(path, dpi=250, bbox_inches='tight', transparent=True, pad_inches=0.05)
    plt.close(fig)

    img = Image.open(path)
    w_px, h_px = img.size
    w_in, h_in = w_px / 250, h_px / 250
    if w_in > max_width_in:
        scale = max_width_in / w_in
        w_in *= scale; h_in *= scale
    return path, Inches(w_in), Inches(h_in)


# ---------- slide helpers ----------
def add_title_slide(title, subtitle, name="", date=""):
    slide = prs.slides.add_slide(title_layout)
    for ph in slide.placeholders:
        i = ph.placeholder_format.idx
        if i == 0:   ph.text = title
        elif i == 1: ph.text = subtitle
        elif i == 12: ph.text = name
        elif i == 13: ph.text = date

def add_section(title, description=""):
    slide = prs.slides.add_slide(section_layout)
    for ph in slide.placeholders:
        i = ph.placeholder_format.idx
        if i == 0: ph.text = title
        elif i == 1: ph.text = description

def add_slide(title, bullets, equations=None):
    """Content slide using 'Title and Text' layout.
    bullets: list of str or (indent_level, str).
    equations: list of (latex_str, left_inches, top_inches[, fontsize]).
    """
    slide = prs.slides.add_slide(content_layout)
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = title
        elif idx == 10:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            for j, item in enumerate(bullets):
                level, text = item if isinstance(item, tuple) else (0, item)
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = Pt(16) if level == 0 else Pt(14)
                p.space_after = Pt(4)
    if equations:
        for eq_spec in equations:
            latex, left, top = eq_spec[0], eq_spec[1], eq_spec[2]
            fs = eq_spec[3] if len(eq_spec) > 3 else 20
            mw = eq_spec[4] if len(eq_spec) > 4 else 7.0
            path, w, h = render_eq(latex, fontsize=fs, max_width_in=mw)
            slide.shapes.add_picture(path, Inches(left), Inches(top), w, h)
    return slide


# ======================================================================
# SLIDES
# ======================================================================

# 1 — Title
add_title_slide(
    "Learning Equilibrium Polymer Distributions",
    "with Continuous Normalizing Flows\nand Diffusion Models",
    "18.337 Final Project", ""
)

# 2 — Motivation
add_section("Motivation",
            "From static snapshots to molecular dynamics")

# 3
add_slide("Motivation", [
    "Goal: infer molecular conformational dynamics from static structural samples",
    "Experimental methods (NMR, CryoEM) provide structural ensembles but not time ordering",
    "Start with controlled Brownian-dynamics synthetic systems",
    "Test whether CNFs and diffusion models can learn equilibrium distributions and score fields",
])

# 4 — Brownian Dynamics
add_slide("Brownian\nDynamics", [
    "N beads in d dimensions",
    "Rouse chain: harmonic springs between adjacent beads",
    "At equilibrium, drift equals score scaled by D",
    "Score target F(x)/D provides a direct test for learned models",
], equations=[
    (r"dX_t = F(X_t)\,dt + \sqrt{2D}\,dW_t", 4.8, 1.3),
    (r"F_i(x) = k_\xi\!\left(x_{i+1} - 2x_i + x_{i-1}\right)", 4.8, 2.8),
    (r"F(x) = D\,\nabla_x \log p_{\mathrm{eq}}(x)", 4.8, 4.3),
])

# 5 — Hamiltonian
add_section("Hamiltonian & Boltzmann Background",
            "Bridging generative modeling and trajectory inference")

# 6
add_slide("Hamiltonian\n& Boltzmann", [
    "Hamilton's equations of motion",
    "Kinetic + potential energy decomposition",
    "Boltzmann distribution at equilibrium",
    "Log-density gradient gives the force",
], equations=[
    (r"\frac{dq_{i,k}}{dt} = \frac{\partial H}{\partial p_{i,k}}"
     r", \quad \frac{dp_{i,k}}{dt} = -\frac{\partial H}{\partial q_{i,k}}",
     4.8, 1.3),
    (r"H = \sum_i \frac{p_{i,k}^2}{2m_i} + V(\vec{q})",
     4.8, 2.5),
    (r"P(\vec{q}) = \frac{1}{Z}\exp\!\left(-\frac{V(\vec{q})}{k_B T}\right)",
     4.8, 3.5),
    (r"\frac{dp_{i,k}}{dt} = k_B T \,\frac{\partial}{\partial q_{i,k}}\log P(\vec{q})",
     4.8, 4.7),
])

# 7 — Data
add_slide("Training\nData", [
    "Centered bead configs from Rouse trajectories",
    "HDF5 arrays, B = number of frames",
    "2D chains, N=32 beads",
    "D=0.25, unit bond length, k_xi=3.0",
    "Center each frame to learn internal conformations",
], equations=[
    (r"x \in \mathbb{R}^{d \times N \times B}", 4.8, 1.7),
])

# 8 — Synthetic Systems
add_section("Synthetic Systems",
            "Base Rouse chain and hairpin-stabilized chain")

# 9 — Base Rouse
add_slide("Base Rouse\nChain", [
    "Quadratic bond potential",
    "Score = negative gradient of potential",
    "N=32, 2D, D=0.25, k_xi=3.0",
    "No nonideal interactions — clean baseline",
], equations=[
    (r"U_{\mathrm{bond}}(x) = \frac{k_\xi}{2D}\sum_{i=2}^{N}\|x_i - x_{i-1}\|^2",
     4.8, 1.3),
    (r"\nabla_x \log p_{\mathrm{eq}}(x) = -\nabla_x U_{\mathrm{bond}}(x)",
     4.8, 3.2),
    (r"F(x) = D\,\nabla_x \log p_{\mathrm{eq}}(x)",
     4.8, 4.8),
])

# 10 — Hairpin
add_slide("Hairpin\nChain", [
    "Nonlocal contacts + excluded volume",
    "LJ interactions with cutoff & softening",
    "Attractive LJ between paired beads",
    "Confinement toward center of mass",
], equations=[
    (r"U_{\mathrm{LJ}}(r) = 4\epsilon\!\left[\!\left(\frac{\sigma}{r}\right)^{\!12} "
     r"- \left(\frac{\sigma}{r}\right)^{\!6}\right]", 4.8, 1.3),
    (r"U_{\mathrm{EV}}(r) = \epsilon_{\mathrm{EV}}\!\left(\frac{\sigma_{\mathrm{EV}}}{r}\right)^p",
     4.8, 3.0),
    (r"U_{\mathrm{conf}}(x) = c\sum_{i=1}^{N}\|x_i - \bar{x}\|^2, \quad c=0.02",
     4.8, 4.5),
])

# 11 — CNF
add_section("Continuous Normalizing Flow",
            "Learning densities via neural ODEs")

# 12 — CNF core
add_slide("CNF\nModel", [
    "ODE maps data to standard normal",
    "tau in [0,1]: data to noise",
    "Log-density via change of variables",
    "Minimize negative log-likelihood",
], equations=[
    (r"\frac{dz(\tau)}{d\tau} = f_\theta\!\left(z(\tau), \tau\right)",
     4.8, 1.3),
    (r"\frac{d}{d\tau}\log p(z(\tau)) = "
     r"-\operatorname{Tr}\!\left(\frac{\partial f_\theta}{\partial z}\right)",
     4.8, 2.8),
    (r"\log p_\theta(x) = \log\mathcal{N}(z(1);\,0,I) - \Delta\!\log p",
     4.8, 4.3),
])

# 13 — Augmented ODE
add_slide("CNF\nAugmented\nODE", [
    "Joint state: coordinates +\naccumulated log-density change",
    "Single ODE solve gives both\ntransformed point and likelihood",
], equations=[
    (r"\begin{bmatrix} z(1) \\ \Delta\!\log P_z(1) \end{bmatrix}"
     r" = \begin{bmatrix} x \\ 0 \end{bmatrix}"
     r" + \int_0^1\!\begin{bmatrix} f_\theta(z,\tau) \\"
     r" -\operatorname{Tr}\!\left(\frac{\partial f_\theta}{\partial z}\right)"
     r" \end{bmatrix} d\tau",
     4.8, 2.5, 18),
    (r"\mathcal{L}_{\mathrm{CNF}}(\theta) = "
     r"-\mathbb{E}_{x \sim p_{\mathrm{data}}}\!\left[\log p_\theta(x)\right]",
     4.8, 5.2),
])

# 14 — Hutchinson
add_slide("Hutchinson\nJVP Trace", [
    "Avoids materializing the full Jacobian",
    "Rademacher probe vectors v",
    "Explicit JVP through EGNN —\none forward pass",
    "Also supports exact trace\nand finite-difference modes",
], equations=[
    (r"\operatorname{Tr}\!\left(\frac{\partial f_\theta}{\partial x}\right)"
     r" = \mathbb{E}_v\!\left[v^\top \frac{\partial f_\theta}{\partial x}\, v\right]",
     4.8, 1.3),
    (r"\widehat{\operatorname{Tr}} = "
     r"\sum_{i,j} v_{i,j}\!\left[J_{f_\theta}(x,\tau)\,v\right]_{i,j}",
     4.8, 3.5),
])

# 15 — CNF Architecture
add_slide("CNF\nArchitecture", [
    "EGNN-style vector field\n(translation equivariant after centering)",
    "Pairwise messages from node\nembeddings, |i-j|, distances, tau",
    "Coordinate updates from\nrelative displacements",
    "SciML ODEProblem + InterpolatingAdjoint\n+ ZygoteVJP for training",
])

# 16 — Diffusion
add_section("Variance-Preserving Diffusion Model",
            "Learning the score field directly")

# 17 — VP forward
add_slide("VP Diffusion\nForward", [
    "Continuous-time VP forward SDE",
    "Linear noise schedule",
    "Marginal distribution at time t",
    "Score target for training",
], equations=[
    (r"dx_t = -\tfrac{1}{2}\beta(t)\,x_t\,dt + \sqrt{\beta(t)}\,dw_t",
     4.8, 1.3),
    (r"\beta(t) = \beta_{\min} + t\,(\beta_{\max} - \beta_{\min})",
     4.8, 2.6),
    (r"x_t = \mu(t)\,x_0 + \sigma(t)\,\epsilon, \quad \epsilon\sim\mathcal{N}(0,I)",
     4.8, 3.8),
    (r"\nabla_{x_t}\!\log p_{0t}(x_t\,|\,x_0) = -\frac{\epsilon}{\sigma(t)}",
     4.8, 5.0),
])

# 18 — Diffusion training & sampling
add_slide("Diffusion\nTraining\n& Sampling", [
    "Weighted denoising score matching",
    "Reverse-time SDE with\nEuler–Maruyama sampler",
    "Stability: clip score norms,\nrecenter, gradient clipping",
], equations=[
    (r"\mathcal{L}_{\mathrm{diff}} = \mathbb{E}_{t,x_0,\epsilon}\!"
     r"\left[\sigma^2(t)\left\|s_\theta(x_t,t) + "
     r"\frac{\epsilon}{\sigma(t)}\right\|^2\right]",
     4.8, 1.3, 18),
    (r"x_{t-\Delta t} = x_t - \left[f(t)\,x_t - g(t)^2\,s_\theta(x_t,t)\right]\Delta t"
     r" + g(t)\sqrt{\Delta t}\,z",
     4.8, 3.8, 16, 7.5),
])

# 19 — Evaluation
add_slide("Evaluation", [
    "MAE of average pairwise distances\n(generated vs. training)",
    "Fraction of finite generated samples",
    "Score-field diagnostics against\nknown Brownian target",
    "CNF: differentiate learned log-likelihood",
    "Diffusion: network output at small t",
], equations=[
    (r"\nabla_x \log p_\theta(x) \;\longleftrightarrow\; \frac{F(x)}{D}",
     4.8, 5.0),
])

# 20 — Julia
add_section("Julia & SciML Implementation",
            "Composable simulation, autodiff,\nand ML in one language")

# 21 — SciML stack
add_slide("SciML\nStack", [
    "Data: SDEProblem + StochasticDiffEq\n(EM, Euler-Heun, SOSRA, ...)",
    "CNF: ODEProblem + ComponentArray\n(coords + log-density)",
    (1, "Adaptive Tsit5; InterpolatingAdjoint"),
    "Both models: Zygote reverse-mode AD",
    "Custom adjoints for batched ops,\nconcatenation, centering",
    "Math maps directly to SciML code",
], equations=[
    (r"dX_t = F(X_t)\,dt + \sqrt{2D}\,dW_t \;\;\longrightarrow\;\; \texttt{SDEProblem}",
     4.8, 5.2, 16),
])

# 22 — Trajectory inference
add_slide("BoltzFlow\nTrajectory\nInference", [
    "1. Obtain equilibrium snapshots",
    "2. Train CNF to model P(q)",
    "3. At each timestep:",
    (1, "Compute log-density gradient"),
    (1, "Boltzmann relation for forces"),
    (1, "Update via Hamilton's equations"),
], equations=[
    (r"p_{i,k}(t_0) \sim \mathcal{N}(0,\; m_i\,k_B T)", 4.8, 1.5),
    (r"-\nabla_{\vec{q}}\,V(\vec{q}) = k_B T\,\nabla_{\vec{q}}\log P(\vec{q})",
     4.8, 3.0),
    (r"\frac{d\vec{p}}{dt} = k_B T\,\nabla_{\vec{q}}\log P(\vec{q})",
     4.8, 4.5),
])

# 23 — Summary
add_slide("Summary", [
    "Brownian-dynamics and polymer-\nequilibrium study",
    "Equivariant CNF and diffusion models\non Rouse equilibrium snapshots",
    "Models generate plausible conformations\nand recover equilibrium scores",
    "Protein dynamics: longer-term goal",
    "Julia + SciML: simulation, ODE solving,\nadjoint AD, GPU — one language",
])

# ======================================================================
prs.save(OUTPUT)
print(f"Saved {OUTPUT} with {len(prs.slides)} slides")
print(f"Equation images in {EQ_DIR}")
